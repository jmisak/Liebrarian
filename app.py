import os
import yaml
import logging
import pandas as pd
import docx
from pptx import Presentation
import gradio as gr
from openai import OpenAI
from sentence_transformers import SentenceTransformer
import chromadb
from chromadb.utils import embedding_functions

# ===================================================================
#                         1. CONFIGURATION AND LOGGING
# ===================================================================

def load_config(config_path='config.yaml'):
    """Loads the YAML configuration file."""
    with open(config_path, 'r') as f:
        return yaml.safe_load(f)

CONFIG = load_config()

# Setup robust logging
logging.basicConfig(
    level=CONFIG['log_level'],
    format='%(asctime)s - %(levelname)s - %(message)s',
    handlers=[
        logging.FileHandler(CONFIG['log_file']),
        logging.StreamHandler()
    ]
)

# ===================================================================
#                         2. FILE PARSING MODULE
# ===================================================================

class FileParser:
    """A modular class to extract text from various file types."""
    
    def extract_text(self, filepath):
        """Main method to delegate parsing based on file extension."""
        _, extension = os.path.splitext(filepath)
        if extension == '.docx':
            return self._extract_from_docx(filepath)
        elif extension == '.pptx':
            return self._extract_from_pptx(filepath)
        elif extension == '.xlsx':
            return self._extract_from_xlsx(filepath)
        else:
            logging.warning(f"Unsupported file type: {filepath}. Skipping.")
            return None

    def _extract_from_docx(self, filepath):
        """Extracts text from a .docx file."""
        try:
            doc = docx.Document(filepath)
            return "\n".join([para.text for para in doc.paragraphs if para.text])
        except Exception as e:
            logging.error(f"Error reading docx file {filepath}: {e}")
            return None

    def _extract_from_pptx(self, filepath):
        """Extracts text from a .pptx file."""
        try:
            pres = Presentation(filepath)
            text_runs = []
            for slide in pres.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text_runs.append(shape.text)
            return "\n".join(text_runs)
        except Exception as e:
            logging.error(f"Error reading pptx file {filepath}: {e}")
            return None

    def _extract_from_xlsx(self, filepath):
        """Extracts text from an .xlsx file, converting sheets to string format."""
        try:
            xls = pd.ExcelFile(filepath)
            content = []
            for sheet_name in xls.sheet_names:
                df = pd.read_excel(xls, sheet_name=sheet_name)
                # Filter out empty rows/columns before converting to string
                df.dropna(how='all', inplace=True)
                df.dropna(how='all', axis=1, inplace=True)
                if not df.empty:
                    content.append(f"Sheet: {sheet_name}\n{df.to_string(index=False)}\n")
            return "\n".join(content)
        except Exception as e:
            logging.error(f"Error reading xlsx file {filepath}: {e}")
            return None

# ===================================================================
#      3. VECTOR STORE MANAGEMENT MODULE (REVISED)
# ===================================================================

class VectorStoreManager:
    """Handles the ChromaDB vector store for persistent storage."""
    
    def __init__(self, config):
        self.config = config
        self.client = chromadb.PersistentClient(path=self.config['vector_store_directory'])
        self.embedding_function = embedding_functions.SentenceTransformerEmbeddingFunction(
            model_name=self.config['embedding_model_name']
        )
        # Check if the collection exists to avoid dimension mismatch errors
        try:
            # Try to get the collection, but explicitly pass the embedding function
            # This is crucial for consistency.
            self.collection = self.client.get_collection(
                name=self.config['collection_name'],
                embedding_function=self.embedding_function
            )
            logging.info("Existing collection found.")
        except Exception as e:
            # If the collection doesn't exist or there's an issue, recreate it.
            logging.warning(f"Collection not found or incompatible: {e}. Recreating...")
            self.collection = self.client.get_or_create_collection(
                name=self.config['collection_name'],
                embedding_function=self.embedding_function
            )

    def add_documents(self, documents, metadatas, ids):
        """Adds or updates documents in the collection."""
        if not documents:
            logging.warning("No documents provided to add to vector store.")
            return
        logging.info(f"Adding {len(documents)} document chunks to the collection.")
        # We need to handle potential add errors, particularly duplicates or malformed data.
        try:
            self.collection.add(documents=documents, metadatas=metadatas, ids=ids)
        except Exception as e:
            logging.critical(f"Failed to add documents to ChromaDB: {e}. The collection may be corrupted or the embedding dimensionality is mismatched. You may need to delete the collection folder manually.")
            # For demonstration, you could add code here to try and delete and recreate the collection.

    def query(self, query_text, k_value, file_type):
        """Queries the collection for relevant documents."""
        logging.info(f"Querying vector store for: '{query_text}'")
        results = self.collection.query(
            query_texts=[query_text],
            n_results=k_value,
            include=['documents']
        )
        return results['documents'][0] if results['documents'] else []

# ===================================================================
#      4. LLM INTERACTION MODULE
# ===================================================================

class LLMClient:
    """A client to interact with the LLM server (e.g., LM Studio)."""
    
    def __init__(self, config):
        self.config = config
        self.client = OpenAI(base_url=config['llm_base_url'], api_key=config['llm_api_key'])

    def get_response(self, query, context):
        """Generates a response from the LLM based on a query and context."""
        system_prompt = (
            "You are an expert assistant. Your task is to answer the user's question based "
            "solely on the provided context. If the information is not in the context, "
            "clearly state that you cannot answer based on the available information."
        )
        
        user_prompt = f"Context:\n---\n{''.join(context)}\n---\nQuestion: {query}"
        
        try:
            response = self.client.chat.completions.create(
                model="local-model", # Corrected line: Use a simple, recognized name
                messages=[
                    {"role": "system", "content": system_prompt},
                    {"role": "user", "content": user_prompt}
                ],
                temperature=0.7,
            )
            return response.choices[0].message.content
        except Exception as e:
            logging.critical(f"Failed to get response from LLM: {e}")
            return "Error: Could not connect to the Language Model. Please ensure LM Studio is running and the server is active."

# ===================================================================
# 5. CORE DIGITAL LIBRARIAN
# ===================================================================

class DigitalLibrarian:
    """The main class orchestrating the entire RAG pipeline."""

    def __init__(self, config):
        self.config = config
        self.parser = FileParser()
        self.vector_store = VectorStoreManager(config)
        self.llm_client = LLMClient(config)

    def index_source_documents(self):
        """Parses and indexes all documents from the source directory."""
        logging.info("Starting document indexing process...")
        source_dir = self.config['source_directory']
        if not os.path.exists(source_dir):
            logging.error(f"Source directory '{source_dir}' not found. Please create it and add documents.")
            return

        documents = []
        metadatas = []
        ids = []
        
        for filename in os.listdir(source_dir):
            filepath = os.path.join(source_dir, filename)
            if os.path.isfile(filepath):
                logging.info(f"Processing file: {filename}")
                text = self.parser.extract_text(filepath)
                if text:
                    chunks = [text[i:i+1000] for i in range(0, len(text), 800)]
                    for i, chunk in enumerate(chunks):
                        doc_id = f"{filename}_{i}"
                        
                        documents.append(f"Source: {filename}\n\n{chunk}")
                        
                        metadatas.append({"source": filename})
                        ids.append(doc_id)

        if documents:
            batch_size = 5000
            for i in range(0, len(documents), batch_size):
                batch_documents = documents[i:i + batch_size]
                batch_metadatas = metadatas[i:i + batch_size]
                batch_ids = ids[i:i + batch_size]
                
                logging.info(f"Adding batch {i//batch_size + 1} of {len(batch_documents)} chunks...")
                self.vector_store.add_documents(batch_documents, batch_metadatas, batch_ids)
            
            logging.info("Document indexing process completed successfully.")
        else:
            logging.warning("No new documents were found or processed during indexing.")

    def answer_query(self, query, use_keyword_search, file_type, k_value):
        """Answers a user query using the RAG pipeline with new UI options."""
        if not query:
            return "Please enter a question."

        if use_keyword_search:
            return self._search_by_keyword(query, file_type)
        else:
            context = self.vector_store.query(query, k_value, file_type)
            logging.info(f"Retrieved context: {context}")
            
            if not context:
                return "I could not find any relevant information in the documents to answer your question."
            
            return self.llm_client.get_response(query, context)

    def _search_by_keyword(self, query, file_type):
        """Performs a simple keyword search across source documents."""
        logging.info(f"Performing keyword search for: '{query}'")
        found_files = []
        source_dir = self.config['source_directory']

        for filename in os.listdir(source_dir):
            if file_type != "All" and not filename.endswith(file_type):
                continue
            
            filepath = os.path.join(source_dir, filename)
            text = self.parser.extract_text(filepath)
            
            if text and query.lower() in text.lower():
                found_files.append(filename)
                
        if found_files:
            return f"Found the following files containing the keyword '{query}':\n\n- " + "\n- ".join(found_files)
        else:
            return "No files found containing that keyword."
# ===================================================================
#      6. GRADIO WEB INTERFACE
# ===================================================================

def create_gradio_interface(librarian):
    """Creates and configures the Gradio web UI."""
    logging.info("Creating Gradio interface.")
    
    with gr.Blocks(theme=gr.themes.Soft()) as interface:
        gr.Markdown(
            """
            # 📚 Liebrarian: A Digital Librarian
            Search the company's files and get intelligent results.
            """
        )
        
        with gr.Row():
            query_input = gr.Textbox(
                lines=3,
                placeholder="e.g., Have we done any studies on X drug for Y company?",
                label="Your Question"
            )
            
        with gr.Row():
            # New components for advanced search
            search_type_checkbox = gr.Checkbox(label="Enable Keyword Search (Bypasses LLM)")
            file_type_dropdown = gr.Dropdown(
                choices=["All", ".docx", ".pptx", ".xlsx"],
                value="All",
                label="Filter by File Type"
            )
            k_value_slider = gr.Slider(
                minimum=1,
                maximum=10,
                value=3,
                step=1,
                label="Number of Retrieved Documents (k-value)"
            )
        
        submit_button = gr.Button("Get Answer", variant="primary")
        
        with gr.Blocks():
            answer_output = gr.Markdown(label="Answer")

        submit_button.click(
            fn=librarian.answer_query,
            inputs=[query_input, search_type_checkbox, file_type_dropdown, k_value_slider],
            outputs=answer_output
        )
        
    return interface

# ===================================================================
#        7. MAIN EXECUTION
# ===================================================================

if __name__ == "__main__":
    # --- Step 1: Initialize the Librarian ---
    librarian = DigitalLibrarian(CONFIG)

    # --- Step 2: Conditionally run the indexing process ---
    # Check if the vector store exists and has documents.
    if librarian.vector_store.collection.count() == 0:
        logging.info("Vector store is empty. Starting document indexing...")
        librarian.index_source_documents()
    else:
        logging.info("Vector store already contains documents. Skipping indexing.")

    # --- Step 3: Launch the Web UI ---
    ui = create_gradio_interface(librarian)

    auth_creds = [tuple(cred) for cred in CONFIG['auth_credentials']]

    logging.info("Launching Gradio app...")
    ui.launch(
        share=CONFIG['share_ui'],
        server_name=CONFIG['server_name'],
        auth=auth_creds
    )