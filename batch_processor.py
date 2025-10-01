import os
import json
import re
from pptx import Presentation
from collections import defaultdict

def extract_keywords(text):
    keywords = []
    for word in ["malaria", "study", "Druga", "Drugb", "HCP", "tracking", "TPP", "Zenket"]:
        if word.lower() in text.lower():
            keywords.append(word)
    return keywords

def extract_project_number(text):
    match = re.search(r"(Project\s+\d+|Study\s+ID:\s+\w+-\d+)", text, re.IGNORECASE)
    return match.group(0) if match else None

def extract_slide_content(slide):
    text = []
    for shape in slide.shapes:
        if hasattr(shape, "text") and shape.text.strip():
            text.append(shape.text.strip())
    if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
        notes = slide.notes_slide.notes_text_frame.text
        if notes:
            text.append(f"Notes: {notes.strip()}")
    return "\n".join(text)

def process_pptx(file_path):
    try:
        prs = Presentation(file_path)
        chunks = []
        for i, slide in enumerate(prs.slides):
            slide_text = extract_slide_content(slide)
            if slide_text:
                project_id = extract_project_number(slide_text)
                chunk = {
                    "text": f"Slide {i+1}: {slide_text}",
                    "metadata": {
                        "slide_number": i + 1,
                        "source_file": os.path.basename(file_path),
                        "project_id": project_id,
                        "keywords": extract_keywords(slide_text)
                    }
                }
                chunks.append(chunk)
        print(f"✅ Processed {len(chunks)} chunks from {os.path.basename(file_path)}")
        return chunks
    except Exception as e:
        print(f"❌ Error processing {file_path}: {e}")
        return []

def batch_process(folder_path, output_jsonl, output_grouped_json):
    all_chunks = []
    grouped_by_project = defaultdict(list)

    for root, dirs, files in os.walk(folder_path):
        for filename in files:
            if filename.lower().endswith(".pptx"):
                file_path = os.path.join(root, filename)
                chunks = process_pptx(file_path)
                all_chunks.extend(chunks)
                for chunk in chunks:
                    pid = chunk["metadata"].get("project_id") or "Unassigned"
                    grouped_by_project[pid].append(chunk)

    if all_chunks:
        os.makedirs(os.path.dirname(output_jsonl), exist_ok=True)

        # Write flat .jsonl file
        with open(output_jsonl, "w", encoding="utf-8") as f:
            for chunk in all_chunks:
                f.write(json.dumps(chunk) + "\n")

        # Write grouped .json file
        with open(output_grouped_json, "w", encoding="utf-8") as f:
            json.dump(grouped_by_project, f, indent=2)

        print(f"\n🎉 Done! {len(all_chunks)} chunks written to {output_jsonl}")
        print(f"📁 Grouped by project ID in {output_grouped_json}")
    else:
        print("\n⚠️ No chunks extracted. Check your files or folder path.")

# Example usage
if __name__ == "__main__":
    input_folder = "G:/Apps/Liebrarian"
    output_jsonl = "G:/Apps/Liebrarian/output.jsonl"
    output_grouped_json = "G:/Apps/Liebrarian/project_summary.json"
    batch_process(input_folder, output_jsonl, output_grouped_json)